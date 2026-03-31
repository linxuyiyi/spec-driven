#!/usr/bin/env python3
# -*- coding: utf-8 -*-

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
DARK_BG = "0F172A"
LIGHT_TEXT = "E2E8F0"
ACCENT_BLUE = "38BDF8"
SUBTEXT = "94A3B8"
WARNING_RED = "F87171"

def hex_to_rgb(h):
    h = h.lstrip('#')
    return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))

def set_fill(shape, hex_color):
    shape.fill.solid()
    r, g, b = hex_to_rgb(hex_color)
    shape.fill.fore_color.rgb = RGBColor(r, g, b)

def set_line(shape, hex_color):
    r, g, b = hex_to_rgb(hex_color)
    shape.line.color.rgb = RGBColor(r, g, b)

def add_text(slide, left, top, width, height, text, size=10, bold=False, color=LIGHT_TEXT, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    r, g, b = hex_to_rgb(color)
    p.font.color.rgb = RGBColor(r, g, b)
    p.alignment = align
    return txBox

def section_title(slide, left, top, text):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(0.06), Inches(0.22))
    set_fill(bar, ACCENT_BLUE)
    bar.line.fill.background()
    add_text(slide, left + 0.12, top, 3, 0.25, text, size=10, bold=True, color=ACCENT_BLUE)

# Create slide
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Background
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
set_fill(bg, DARK_BG)
bg.line.fill.background()

# Title
add_text(slide, 0.4, 0.2, 12, 0.5, "范式类测试 E2E 测试方案", size=26, bold=True, color="FFFFFF", align=PP_ALIGN.CENTER)

# Subtitle / 目标
add_text(slide, 0.4, 0.7, 12, 0.3, "目标：从需求到测试结果端到端测试", size=12, color=SUBTEXT, align=PP_ALIGN.CENTER)

# Tags
for i, tag in enumerate(["参数约束", "参数升级", "License 校验"]):
    x = 4.8 + i * 1.3
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.0), Inches(1.1), Inches(0.3))
    set_fill(s, "1E3A5F")
    set_line(s, ACCENT_BLUE)
    add_text(slide, x + 0.05, 1.02, 1, 0.25, tag, size=9, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)

# ===== 业务流程 (left column, top) =====
section_title(slide, 0.4, 1.5, "业务流程")

biz_steps = [("需求描述", "原始需求"), ("文本用例", "结构化要素"), ("脚本用例", "可执行脚本"), ("测试结果", "E2E反馈")]
for i, (t, s) in enumerate(biz_steps):
    x = 0.5 + i * 2.9
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.8), Inches(2.5), Inches(0.7))
    set_fill(box, "1E3A5F")
    set_line(box, ACCENT_BLUE if i == 3 else "334455")
    add_text(slide, x + 0.1, 1.88, 2.3, 0.3, t, size=11, bold=True, color=LIGHT_TEXT, align=PP_ALIGN.CENTER)
    add_text(slide, x + 0.1, 2.15, 2.3, 0.25, s, size=8, color=SUBTEXT, align=PP_ALIGN.CENTER)
    if i < 3:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + 2.55), Inches(2.05), Inches(0.3), Inches(0.12))
        set_fill(arrow, ACCENT_BLUE)
        arrow.line.fill.background()

# ===== E2E动作流程 (full width, below) =====
section_title(slide, 0.4, 2.7, "E2E 动作流程")

# Row 1
e2e1 = [("获取原始需求", "Input"), ("结构化测试要素", "TFO"), ("设计文本用例", "TFO"),
        ("审核文本用例", "MDE"), ("生成脚本用例", "AE"), ("匹配测试环境", "AE")]
for i, (t, a) in enumerate(e2e1):
    x = 0.4 + i * 2.08
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(3.0), Inches(1.95), Inches(0.55))
    set_fill(box, "1E3A5F")
    set_line(box, ACCENT_BLUE if a == "MDE" else "334455")
    add_text(slide, x + 0.08, 3.04, 1.8, 0.25, t, size=8, bold=True, color=LIGHT_TEXT, align=PP_ALIGN.CENTER)
    add_text(slide, x + 0.08, 3.28, 1.8, 0.2, a, size=7, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)
    if i < 5:
        arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + 2.0), Inches(3.2), Inches(0.06), Inches(0.1))
        set_fill(arr, "556677")
        arr.line.fill.background()

# Row 2
e2e2 = [("运行调试脚本", "AE"), ("创建上库MR", "AE"), ("审核脚本用例", "MDE"), ("创建流水线任务", "TE")]
for i, (t, a) in enumerate(e2e2):
    x = 2.5 + i * 2.6
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(3.65), Inches(2.35), Inches(0.55))
    set_fill(box, "1E3A5F")
    set_line(box, ACCENT_BLUE if a in ["MDE", "TE"] else "334455")
    add_text(slide, x + 0.08, 3.69, 2.2, 0.25, t, size=8, bold=True, color=LIGHT_TEXT, align=PP_ALIGN.CENTER)
    add_text(slide, x + 0.08, 3.93, 2.2, 0.2, a, size=7, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)
    if i < 3:
        arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + 2.4), Inches(3.85), Inches(0.15), Inches(0.1))
        set_fill(arr, "556677")
        arr.line.fill.background()

# Connector arrow
conn = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(12.3), Inches(3.55), Inches(0.12), Inches(0.1))
set_fill(conn, "556677")
conn.line.fill.background()

# ===== 难点 (left column) =====
section_title(slide, 0.4, 4.4, "难点与挑战")

challenges = [
    "AI 缺少无线业务背景，业务规则或惯例可能遗漏",
    "AI 理解原始需求偏差导致整体方向错误",
    "AI 调试时对环境/产品问题识别困难"
]
for i, ch in enumerate(challenges):
    y = 4.7 + i * 0.42
    # Warning icon
    w = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(y), Inches(0.25), Inches(0.25))
    set_fill(w, "7F1D1D")
    w.line.fill.background()
    add_text(slide, 0.44, y + 0.02, 0.2, 0.2, "!", size=10, bold=True, color=WARNING_RED, align=PP_ALIGN.CENTER)
    add_text(slide, 0.75, y, 5.5, 0.35, ch, size=9, color=LIGHT_TEXT)

# ===== 解决措施 (right column) =====
section_title(slide, 6.8, 4.4, "解决措施")

# 措施1
n1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.8), Inches(4.75), Inches(0.25), Inches(0.25))
set_fill(n1, ACCENT_BLUE)
n1.line.fill.background()
add_text(slide, 6.83, 4.77, 0.2, 0.2, "1", size=9, bold=True, color="FFFFFF", align=PP_ALIGN.CENTER)
add_text(slide, 7.15, 4.75, 5.5, 0.25, "构建范式测试指导书", size=10, bold=True, color=LIGHT_TEXT)

items1 = ["参数 OCL 设计指导书", "参数升级设计指导书", "无线参数定义说明书", "无线 License 定义说明书"]
for i, item in enumerate(items1):
    add_text(slide, 7.3, 5.05 + i * 0.25, 5, 0.22, "• " + item, size=8, color=SUBTEXT)

# 措施2
n2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.8), Inches(6.1), Inches(0.25), Inches(0.25))
set_fill(n2, ACCENT_BLUE)
n2.line.fill.background()
add_text(slide, 6.83, 6.12, 0.2, 0.2, "2", size=9, bold=True, color="FFFFFF", align=PP_ALIGN.CENTER)
add_text(slide, 7.15, 6.1, 5.5, 0.25, "多 Agent 协作智能团队", size=10, bold=True, color=LIGHT_TEXT)

# Agent mini cards
agents = [("MDE", "测试设计审核人"), ("TFO", "测试设计负责人"), ("AE", "脚本工程师"), ("TE", "测试执行人")]
for i, (abbr, role) in enumerate(agents):
    x = 7.3 + (i % 2) * 3.0
    y = 6.4 + (i // 2) * 0.35
    add_text(slide, x, y, 0.5, 0.25, abbr, size=8, bold=True, color=ACCENT_BLUE)
    add_text(slide, x + 0.5, y, 2.2, 0.25, role, size=8, color=SUBTEXT)

# ===== Bottom badges =====
badges = ["E2E 端到端", "多 Agent 协作", "范式测试驱动"]
for i, b in enumerate(badges):
    x = 4.5 + i * 1.8
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(7.05), Inches(1.5), Inches(0.32))
    set_fill(s, "1E3A5F")
    set_line(s, ACCENT_BLUE)
    add_text(slide, x + 0.05, 7.07, 1.4, 0.28, b, size=9, bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)

# Save
output_path = "/Users/niuniu/程序/testagent/范式类测试E2E测试方案.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")